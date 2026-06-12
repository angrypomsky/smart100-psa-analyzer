from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import pptx.util as util

# Color scheme
C_BLUE = RGBColor(0x1F, 0x4E, 0x79)   # dark blue (title bg)
C_LBLUE = RGBColor(0x2E, 0x75, 0xB6)  # mid blue (accent)
C_ACCENT = RGBColor(0xED, 0x7D, 0x31) # orange accent
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xD6, 0xE4, 0xF0)  # light blue bg
C_DARK = RGBColor(0x1A, 0x1A, 0x2E)   # near black text
C_GRAY = RGBColor(0x70, 0x70, 0x70)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # blank layout

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        if line_width:
            shape.line.width = util.Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_bullet_box(slide, title, bullets, l, t, w, h,
                   title_size=16, bullet_size=13,
                   box_fill=C_LIGHT, title_color=C_LBLUE, bullet_color=C_DARK):
    # box background
    add_rect(slide, l, t, w, h, fill_rgb=box_fill)
    # title
    add_text(slide, title, l+0.15, t+0.1, w-0.3, 0.4,
             size=title_size, bold=True, color=title_color)
    # bullets
    txBox = slide.shapes.add_textbox(Inches(l+0.15), Inches(t+0.55),
                                     Inches(w-0.3), Inches(h-0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(bullet_size)
        p.font.color.rgb = bullet_color
        p.space_after = Pt(4)

# ─────────────────────────────────────────────
# Slide 1: Title
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=C_BLUE)
add_rect(sl, 0, 5.6, 13.33, 1.9, fill_rgb=C_LBLUE)

add_text(sl, "SMART PSA를 위한 MARS-KS 출력 자동 후처리 도구",
         0.6, 1.0, 12.1, 1.1, size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "열수력 결과에서 대화형 DET 대시보드까지",
         0.6, 2.2, 12.1, 0.7, size=22, bold=False, color=C_LIGHT, align=PP_ALIGN.CENTER)

add_rect(sl, 5.5, 3.1, 2.3, 0.05, fill_rgb=C_ACCENT)

add_text(sl, "경희대학교 원자력공학과",
         0.6, 5.75, 12.1, 0.5, size=16, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "ANS Winter Meeting 2025",
         0.6, 6.3, 12.1, 0.5, size=14, color=C_LIGHT, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────────────────────────────
# Slide 2: 문제 배경 & 연구 동기
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill_rgb=C_BLUE)
add_text(sl, "배경 및 연구 동기", 0.4, 0.15, 12.5, 0.8,
         size=24, bold=True, color=C_WHITE)

add_bullet_box(sl, "현행 PSA 실무의 한계",
    ["수백 개 TH 시뮬레이션 시나리오 → 수작업 후처리",
     "전문 인력의 반복 투입 → 시간·비용 선형 증가",
     "피로 유발 오류 → 시나리오 간 판정 일관성 저하",
     "이진 사건수목: 부분적 계통 성능 정보 소실"],
    0.4, 1.3, 6.0, 4.0)

add_bullet_box(sl, "SMART 원자로 특성",
    ["100 MWth 일체형 가압경수로 (KAERI 개발)",
     "4계통 피동잔열제거계통(PRHRS) 보유",
     "작동 계통수(0~4대)가 사고 결과에 직접 영향",
     "이진 ET로는 부분 성능 분석 불가"],
    6.6, 1.3, 6.3, 4.0)

add_rect(sl, 0.4, 5.5, 12.5, 0.8, fill_rgb=C_ACCENT)
add_text(sl, "→  자동화된 후처리 + 동적 사건수목(DET) 구조가 필요",
         0.7, 5.6, 12.0, 0.6, size=15, bold=True, color=C_WHITE)

# ─────────────────────────────────────────────
# Slide 3: 도구 파이프라인 개요
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill_rgb=C_BLUE)
add_text(sl, "도구 파이프라인 개요", 0.4, 0.15, 12.5, 0.8,
         size=24, bold=True, color=C_WHITE)

# pipeline boxes
boxes = [
    ("MARS-KS\n출력 Excel", 0.5),
    ("VarMapper\n변수 자동 매핑", 3.0),
    ("변수 추출\n(RT / PCT / PRHRS)", 5.5),
    ("대화형 Excel\n대시보드", 8.0),
    ("PSA 소프트웨어\n(향후)", 10.5),
]
for i, (label, x) in enumerate(boxes):
    color = C_LBLUE if i < 4 else C_GRAY
    add_rect(sl, x, 2.0, 2.3, 1.4, fill_rgb=color)
    add_text(sl, label, x+0.1, 2.15, 2.1, 1.1,
             size=13, bold=(i<4), color=C_WHITE, align=PP_ALIGN.CENTER)
    if i < len(boxes)-1:
        add_text(sl, "→", x+2.3, 2.45, 0.7, 0.5,
                 size=20, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

add_text(sl, "단일 Python 스크립트 (psa_analyzer.py) — GUI 파일 선택만으로 전 과정 자동 실행",
         0.4, 3.7, 12.5, 0.5, size=14, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

# OOP structure
add_rect(sl, 0.4, 4.4, 12.5, 2.7, fill_rgb=C_LIGHT)
add_text(sl, "객체지향 구조 (상속)", 0.6, 4.5, 4.0, 0.4,
         size=15, bold=True, color=C_LBLUE)

base_x = 0.6
add_rect(sl, base_x, 5.0, 3.0, 0.9, fill_rgb=C_LBLUE)
add_text(sl, "BaseAnalyzer\n(공통 로직: RT·PCT·PRHRS)", base_x+0.1, 5.05, 2.8, 0.8,
         size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

children = ["LOFW", "SBLOCA", "GTRN", "LSSB", "SGTR"]
for i, name in enumerate(children):
    cx = 4.2 + i * 1.75
    add_rect(sl, cx, 5.0, 1.55, 0.9, fill_rgb=C_ACCENT)
    add_text(sl, f"{name}\nAnalyzer", cx+0.05, 5.05, 1.45, 0.8,
             size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "↑  파라미터만 각자 정의 (신규 사고유형 추가 용이)",
         4.0, 6.05, 9.0, 0.4, size=11, color=C_GRAY, italic=True)

# ─────────────────────────────────────────────
# Slide 4: 핵심 알고리즘
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill_rgb=C_BLUE)
add_text(sl, "핵심 추출 알고리즘", 0.4, 0.15, 12.5, 0.8,
         size=24, bold=True, color=C_WHITE)

add_bullet_box(sl, "① 원자로 트립(RT) 감지",
    ["기준: rktpow ≤ 초기값 × 20% 최초 시점",
     "PCT 구간 시작: 더 엄격한 10% 기준 적용",
     "이유: 트립 후 관성 잔류출력 수십 초 지속 → PCT는 냉각 단계에서 추출",
     "아티팩트 처리: 첫 행 rktpow=0 → 첫 번째 양수값으로 초기화"],
    0.4, 1.2, 6.0, 3.0)

add_bullet_box(sl, "② PRHRS 작동 계통수 판정",
    ["상대비교 기준: Qi ≥ Qmax × R기준 (기본 10%)",
     "Floor = 1×10⁶ W (전 사고유형 공통)",
     "초기 조건 변화에 강인한 자동 적응",
     "SBLOCA: RT+대기 후 초기 3h 창 집계 (붕괴열 감소 대응)",
     "SGTR: 전체 평균 (후반 열출력 감소 오분류 방지)"],
    6.6, 1.2, 6.3, 3.0)

add_bullet_box(sl, "③ PCT 추출 및 노심손상(CD) 판정",
    ["10% 기준 RT 이후 구간에서 최대 피복재 온도 추출",
     "10,000 K 초과 비물리 이상값 제거",
     "NRC 10 CFR 50.46 기준: PCT ≥ 1477 K → CD"],
    0.4, 4.4, 12.5, 2.6)

# ─────────────────────────────────────────────
# Slide 5: 사고유형별 파라미터 테이블
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill_rgb=C_BLUE)
add_text(sl, "사고유형별 PRHRS 알고리즘 파라미터", 0.4, 0.15, 12.5, 0.8,
         size=24, bold=True, color=C_WHITE)

headers = ["파라미터", "LOFW", "SBLOCA", "GTRN", "LSSB", "SGTR"]
rows = [
    ["집계 구간", "전체 평균", "초기 3h", "전체 평균", "전체 평균", "전체 평균"],
    ["대기 시간 (s)", "100", "100", "100", "100", "100"],
    ["Floor (W)", "1×10⁶", "1×10⁶", "1×10⁶", "1×10⁶", "1×10⁶"],
    ["상대비교 기준", "10%", "10%", "10%", "10%", "10%"],
]

col_w = [2.5, 1.8, 1.8, 1.8, 1.8, 1.8]
col_x = [0.4]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

row_h = 0.65
start_y = 1.4

# header row
for j, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    add_rect(sl, x, start_y, w-0.05, row_h, fill_rgb=C_LBLUE)
    add_text(sl, h, x+0.05, start_y+0.12, w-0.15, row_h-0.2,
             size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows):
    y = start_y + row_h * (i+1)
    bg = C_LIGHT if i % 2 == 0 else C_WHITE
    for j, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
        add_rect(sl, x, y, w-0.05, row_h, fill_rgb=bg,
                 line_rgb=RGBColor(0xCC,0xCC,0xCC), line_width=0.5)
        c = C_LBLUE if j == 0 else C_DARK
        add_text(sl, cell, x+0.05, y+0.12, w-0.15, row_h-0.2,
                 size=13, bold=(j==0), color=c, align=PP_ALIGN.CENTER)

add_text(sl, "* SBLOCA만 초기 3h 창: RT 이후 붕괴열 자연 감소로 후반 출력 하락 → 전체 평균 시 계통수 과소평가 방지",
         0.4, 5.2, 12.5, 0.5, size=11, color=C_GRAY, italic=True)
add_text(sl, "* Floor: 비작동 계통 잔류출력(수백 kW 이하)과 작동 계통(수 MW)을 명확히 분리",
         0.4, 5.7, 12.5, 0.5, size=11, color=C_GRAY, italic=True)

# ─────────────────────────────────────────────
# Slide 6: 대화형 Excel 대시보드
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill_rgb=C_BLUE)
add_text(sl, "대화형 Excel 대시보드 — DET 구조 구현", 0.4, 0.15, 12.5, 0.8,
         size=24, bold=True, color=C_WHITE)

features = [
    ("PRHRS 계통수 필터", "드롭다운 선택(0~4대)\n→ PCT 분포 차트 즉시 갱신"),
    ("PCT 분포 시각화", "히스토그램 + 통계 요약\n(min / mean / max / P95)"),
    ("시나리오 목록", "파일명, RT 결과, PCT,\nCD 판정 일람"),
    ("전체 분포 비교", "계통수별 PCT 분포 중첩\n→ 단조성 확인"),
]
for i, (title, desc) in enumerate(features):
    x = 0.4 + i * 3.2
    add_rect(sl, x, 1.3, 3.0, 2.5, fill_rgb=C_LBLUE)
    add_text(sl, title, x+0.1, 1.4, 2.8, 0.55,
             size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, desc, x+0.1, 2.05, 2.8, 1.6,
             size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)

add_rect(sl, 0.4, 4.0, 12.5, 0.05, fill_rgb=C_ACCENT)

add_text(sl, "DET와의 동등성", 0.4, 4.2, 4.0, 0.45,
         size=16, bold=True, color=C_LBLUE)
add_bullet_box(sl, "",
    ["이진 성공/실패 → PRHRS 계통수(0~4)를 이산 분기 변수로 정의",
     "성공기준 변경 시나리오 즉시 탐색 (예: '3계통' → '2계통 이상'으로 완화)",
     "추가 시뮬레이션 없이 노심손상 빈도 영향 정량 추정 가능"],
    0.4, 4.55, 12.5, 2.6, title_size=1, bullet_size=14,
    box_fill=C_LIGHT, bullet_color=C_DARK)

# ─────────────────────────────────────────────
# Slide 7: 검증 결과
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill_rgb=C_BLUE)
add_text(sl, "검증 결과", 0.4, 0.15, 12.5, 0.8, size=24, bold=True, color=C_WHITE)

# metric cards
metrics = [
    ("RT 감지\n정확도", "100%", "실제 MARS-KS 데이터\n5개 사고유형 수작업 비교"),
    ("PRHRS 계통수\n판정 정확도", "100%", "실제 MARS-KS 데이터\n5개 사고유형 수작업 비교"),
    ("데모 데이터셋\n물리 규칙 위반", "0건", "1,000 시나리오\n4개 일관성 규칙"),
    ("처리 시간", "< 30초", "1,000 시나리오\n일반 PC 기준"),
]
for i, (label, value, note) in enumerate(metrics):
    x = 0.4 + i * 3.2
    add_rect(sl, x, 1.3, 3.0, 2.8, fill_rgb=C_LIGHT)
    add_rect(sl, x, 1.3, 3.0, 0.6, fill_rgb=C_LBLUE)
    add_text(sl, label, x+0.1, 1.35, 2.8, 0.5,
             size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, value, x+0.1, 2.0, 2.8, 0.9,
             size=30, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text(sl, note, x+0.1, 3.0, 2.8, 0.9,
             size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

add_text(sl, "물리적 일관성 검증 규칙 (데모 데이터셋)", 0.4, 4.35, 12.5, 0.45,
         size=15, bold=True, color=C_LBLUE)
rules = [
    ("규칙 A", "PCT–결과 일관성 (1477 K 임계값)"),
    ("규칙 B", "PRHRS 단조성 — 계통수 증가 → PCT 감소"),
    ("규칙 C", "ATWS 패널티 — RT 실패 → PCT 상승"),
    ("규칙 D", "Feed-and-Bleed 작동 조건 제약"),
]
for i, (tag, desc) in enumerate(rules):
    x = 0.4 + (i % 2) * 6.5
    y = 4.9 + (i // 2) * 0.65
    add_rect(sl, x, y, 1.1, 0.5, fill_rgb=C_ACCENT)
    add_text(sl, tag, x+0.05, y+0.08, 1.0, 0.35,
             size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, desc, x+1.2, y+0.08, 5.1, 0.4,
             size=12, color=C_DARK)

add_text(sl, "※ PCT 추출 검증: 현재 700~1,400 K (OK 시나리오) — CD 시나리오 검증은 향후 추가 시뮬레이션 후 수행 예정",
         0.4, 6.9, 12.5, 0.4, size=10, color=C_GRAY, italic=True)

# ─────────────────────────────────────────────
# Slide 8: 결론 및 향후 과제
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill_rgb=C_BLUE)
add_text(sl, "결론 및 향후 과제", 0.4, 0.15, 12.5, 0.8, size=24, bold=True, color=C_WHITE)

add_bullet_box(sl, "핵심 기여",
    ["MARS-KS TH 출력 → 대화형 Excel 대시보드 완전 자동화",
     "단일 Python 스크립트, GUI 파일 선택만으로 전 과정 처리",
     "상대비교 기반 PRHRS 판정: 초기 조건 변화에 강인",
     "VarMapper: 노형 변경 시 코드 수정 없이 재적용 가능",
     "DET 구조로 성공기준 민감도 즉시 탐색 (추가 시뮬레이션 불필요)"],
    0.4, 1.25, 6.1, 4.2)

add_bullet_box(sl, "향후 과제",
    ["실제 MARS-KS 대용량 데이터셋 적용 및 성능 확인",
     "ADS·PSIS·SIT 등 추가 안전계통 확장",
     "AIMS-PSA 호환 사건수목 형식(.ket) 직접 출력",
     "TH 시뮬레이션 → PSA 소프트웨어 완전 자동 파이프라인"],
    6.7, 1.25, 6.2, 4.2)

add_rect(sl, 0.4, 5.65, 12.5, 1.5, fill_rgb=C_ACCENT)
add_text(sl, "목표: MARS-KS 출력에서 PSA 소프트웨어 입력 파일까지\n완전 자동 생성하는 통합 후처리 체계 구축",
         0.8, 5.75, 11.7, 1.3, size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# save
out = r"C:\Users\user\smart100-psa-analyzer\ANS_2025_presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
